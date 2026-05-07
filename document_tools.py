from __future__ import annotations

import csv
import json
import re
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any

from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor

import workspace
from config import Settings


REPORT_MAX_CHARS = 45000
BINARY_PREVIEW_BYTES = 4096

CODE_AND_TEXT_SUFFIXES = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".scss", ".json",
    ".md", ".txt", ".yaml", ".yml", ".toml", ".ini", ".csv", ".sql", ".sh",
    ".ps1", ".bat", ".java", ".cs", ".cpp", ".c", ".h", ".hpp", ".rs", ".go",
    ".php", ".rb", ".swift", ".kt", ".kts", ".xml", ".tex", ".bib", ".m",
    ".r", ".R", ".scala", ".lua", ".pl", ".dart", ".vue", ".svelte", ".erl",
    ".ex", ".exs", ".jl", ".f90", ".f", ".vb", ".fs", ".fsx", ".clj", ".cljs",
}

LANGUAGE_BY_SUFFIX = {
    ".m": "MATLAB/Octave or Objective-C source",
    ".py": "Python source",
    ".tex": "LaTeX document",
    ".pptx": "PowerPoint presentation",
    ".docx": "Word document",
    ".pdf": "PDF document",
    ".ipynb": "Jupyter notebook",
}



DEFAULT_DOCX_FONT = "Times New Roman"
DOCX_BLACK = RGBColor(0, 0, 0)


def _format_run(run, *, size: int, bold: bool = False, font_name: str = DEFAULT_DOCX_FONT) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = DOCX_BLACK
    # Word stores East Asian font separately; setting it keeps formatting more consistent.
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)


def _format_paragraph(paragraph, kind: str) -> None:
    if kind == "title":
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        size = 28
        bold = True
    elif kind == "heading_1":
        size = 16
        bold = True
    elif kind == "heading_2":
        size = 14
        bold = True
    elif kind == "heading_3":
        size = 12
        bold = True
    else:
        size = 11
        bold = False

    for run in paragraph.runs:
        _format_run(run, size=size, bold=bold)

def _sanitize_xml_text(value: Any) -> str:
    text = str(value or "")
    # XML 1.0 forbids null bytes and most C0 control characters. Notebook outputs can contain them.
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]", "", text)
    text = re.sub(r"[\uD800-\uDFFF]", "", text)
    return text
class DocumentToolError(RuntimeError):
    pass


def _truncate(text: str, max_chars: int) -> str:
    text = _sanitize_xml_text(text)
    return text[:max_chars] + ("\n\n[truncated]" if len(text) > max_chars else "")


def _read_text_any_encoding(path: Path, max_chars: int) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1", "utf-16"):
        try:
            return _truncate(data.decode(encoding, errors="strict"), max_chars)
        except UnicodeDecodeError:
            continue
    return _truncate(data.decode("utf-8", errors="replace"), max_chars)


def extract_pdf_text(path: Path, max_chars: int = REPORT_MAX_CHARS) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentToolError("PDF support requires pypdf. Install dependencies from requirements.txt.") from exc

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append(f"\n\n--- Page {index} ---\n{text.strip()}")
        if sum(len(item) for item in pages) >= max_chars:
            break
    return _truncate("".join(pages).strip(), max_chars)


def extract_docx_text(path: Path, max_chars: int = REPORT_MAX_CHARS) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentToolError("DOCX support requires python-docx. Install dependencies from requirements.txt.") from exc

    doc = Document(str(path))
    parts = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            style = paragraph.style.name if paragraph.style else "Normal"
            parts.append(f"[{style}] {text}")
    for table_index, table in enumerate(doc.tables, start=1):
        parts.append(f"\n[Table {table_index}]")
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            parts.append(" | ".join(cells))
    return _truncate("\n".join(parts), max_chars)


def extract_pptx_text(path: Path, max_chars: int = REPORT_MAX_CHARS) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentToolError("PPTX support requires python-pptx. Install dependencies from requirements.txt.") from exc

    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = [f"--- Slide {index} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        slides.append("\n".join(lines))
    return _truncate("\n\n".join(slides), max_chars)


def extract_notebook_text(path: Path, max_chars: int = REPORT_MAX_CHARS) -> str:
    try:
        raw = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    except json.JSONDecodeError as exc:
        raise DocumentToolError("Could not parse the notebook JSON.") from exc

    cells = raw.get("cells", [])
    sections = [
        "# Notebook Extract",
        f"- Total cells: {len(cells)}",
        f"- Markdown cells: {sum(1 for cell in cells if cell.get('cell_type') == 'markdown')}",
        f"- Code cells: {sum(1 for cell in cells if cell.get('cell_type') == 'code')}",
    ]
    figures = 0
    for index, cell in enumerate(cells, start=1):
        source = "".join(cell.get("source", [])).strip()
        if source:
            cell_type = cell.get("cell_type", "cell")
            fence = "python" if cell_type == "code" else "markdown"
            sections.append(f"\n## {cell_type.title()} Cell {index}\n```{fence}\n{source}\n```")
        for output in cell.get("outputs", []):
            if "data" in output and any(str(key).startswith("image/") for key in output.get("data", {})):
                figures += 1
            text_items = output.get("text") or output.get("data", {}).get("text/plain") or []
            if isinstance(text_items, str):
                text_items = [text_items]
            if text_items:
                sections.append(f"\n### Output From Cell {index}\n" + "".join(text_items).strip())
    sections.insert(4, f"- Embedded/inline figure outputs detected: {figures}")
    return _truncate("\n".join(sections), max_chars)


def extract_spreadsheet_text(path: Path, max_chars: int = REPORT_MAX_CHARS) -> str:
    suffix = path.suffix.lower()
    if suffix == ".csv":
        rows = []
        with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            reader = csv.reader(handle)
            for index, row in enumerate(reader):
                if index >= 60:
                    break
                rows.append(" | ".join(row))
        return _truncate("\n".join(rows), max_chars)

    try:
        import openpyxl
    except ImportError as exc:
        raise DocumentToolError("Spreadsheet support requires openpyxl for .xlsx/.xlsm files.") from exc

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets[:10]:
        parts.append(f"--- Sheet: {sheet.title} ---")
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 60:
                break
            parts.append(" | ".join("" if value is None else str(value) for value in row))
    return _truncate("\n".join(parts), max_chars)


def extract_office_zip_text(path: Path, max_chars: int = REPORT_MAX_CHARS) -> str:
    # Fallback for office XML formats when optional libraries are unavailable.
    chunks = []
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.endswith(".xml"):
                continue
            if not any(part in name for part in ("word/", "ppt/", "xl/")):
                continue
            raw = archive.read(name).decode("utf-8", errors="replace")
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
            if text:
                chunks.append(f"--- {name} ---\n{text}")
            if sum(len(item) for item in chunks) >= max_chars:
                break
    return _truncate("\n\n".join(chunks), max_chars)


def _generic_binary_profile(path: Path) -> str:
    data = path.read_bytes()[:BINARY_PREVIEW_BYTES]
    printable = "".join(chr(byte) if 32 <= byte <= 126 else "." for byte in data)
    return (
        "This file is not a directly supported readable text/document format. "
        "A binary profile is provided for identification only.\n\n"
        f"First {len(data)} bytes as printable preview:\n{printable}"
    )


def _extract_generic_content(path: Path, max_chars: int) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".ipynb":
        return extract_notebook_text(path, max_chars), "jupyter_notebook"
    if suffix == ".pdf":
        return extract_pdf_text(path, max_chars), "pdf"
    if suffix == ".docx":
        return extract_docx_text(path, max_chars), "word_document"
    if suffix == ".pptx":
        try:
            return extract_pptx_text(path, max_chars), "powerpoint_presentation"
        except DocumentToolError:
            return extract_office_zip_text(path, max_chars), "powerpoint_presentation_xml_fallback"
    if suffix in {".xlsx", ".xlsm", ".csv"}:
        return extract_spreadsheet_text(path, max_chars), "spreadsheet"
    if suffix in CODE_AND_TEXT_SUFFIXES or workspace.looks_textual(path):
        return _read_text_any_encoding(path, max_chars), LANGUAGE_BY_SUFFIX.get(suffix, "text_or_code")

    try:
        text = _read_text_any_encoding(path, max_chars)
        if text and text.count("\x00") < 5:
            return text, "unknown_text_like_file"
    except Exception:
        pass
    return _generic_binary_profile(path), "binary_or_unsupported_file"


def _keywords_from_instructions(instructions: str) -> list[str]:
    words = re.findall(r"[A-Za-z_][A-Za-z0-9_+.#-]{2,}", instructions.lower())
    stop = {
        "the", "and", "for", "with", "from", "this", "that", "file", "report", "technical",
        "analyze", "analysis", "create", "should", "contain", "names", "all", "also", "use",
        "format", "normal", "text", "headings", "provided", "prompt", "look", "condition", "check", "include", "includes", "including", "limitation", "limitations", "summarize", "summary", "detail", "detailed", "professional", "formal",
    }
    seen = []
    for word in words:
        if word not in stop and word not in seen:
            seen.append(word)
    return seen[:40]


def _condition_matches(content: str, instructions: str) -> list[dict[str, Any]]:
    matches = []
    lowered = content.lower()
    for keyword in _keywords_from_instructions(instructions):
        count = lowered.count(keyword.lower())
        if count:
            first = lowered.find(keyword.lower())
            start = max(0, first - 120)
            end = min(len(content), first + 220)
            matches.append({"condition_term": keyword, "occurrences": count, "sample_context": content[start:end].replace("\n", " ")})
        else:
            matches.append({"condition_term": keyword, "occurrences": 0, "sample_context": ""})
    return matches


def _generic_code_signals(content: str) -> dict[str, list[str]]:
    patterns = {
        "functions_or_methods": r"(?:def|function|func|sub|void|int|float|double|public|private|protected)\s+([A-Za-z_][A-Za-z0-9_]*)\s*\(",
        "classes": r"(?:class|interface|struct)\s+([A-Za-z_][A-Za-z0-9_]*)",
        "imports_or_dependencies": r"^(?:import|from|require|using|include|library|package)\s+([^\n;]+)",
        "figures_or_plots": r"\b(?:plot|figure|imshow|scatter|bar|hist|savefig|ggplot|chart)\b",
        "models_or_algorithms": r"\b(?:regression|classifier|model|network|lstm|cnn|transformer|randomforest|svm|kmeans|linear|tree)\b",
        "metrics_or_results": r"\b(?:accuracy|precision|recall|f1|mae|mse|rmse|r2|loss|auc|score|result)\b",
        "preprocessing_terms": r"\b(?:normalize|standardize|scale|fillna|dropna|missing|encode|tokenize|clean|filter|resample|interpolate|preprocess)\b",
    }
    signals = {}
    for name, pattern in patterns.items():
        found = re.findall(pattern, content, flags=re.IGNORECASE | re.MULTILINE)
        if found:
            cleaned = []
            for item in found[:50]:
                value = item if isinstance(item, str) else " ".join(item)
                value = value.strip()
                if value and value not in cleaned:
                    cleaned.append(value[:120])
            signals[name] = cleaned
    return signals



def _condition_lines_from_signals(content: str, instructions: str, signals: dict[str, list[str]]) -> list[str]:
    categories = {
        "preprocessing": ["preprocess", "preprocessing", "normalize", "normalized", "standardize", "scale", "clean", "encode", "missing", "fillna", "dropna"],
        "models": ["model", "models", "algorithm", "regression", "classifier", "network", "lstm", "svm", "tree"],
        "metrics": ["metric", "metrics", "accuracy", "precision", "recall", "f1", "mae", "mse", "rmse", "r2", "loss", "score"],
        "results": ["result", "results", "output", "outputs", "evaluation", "performance", "score"],
        "figures": ["figure", "figures", "plot", "plots", "chart", "graph", "image", "visualization"],
        "tables": ["table", "tables", "tabular"],
        "references": ["reference", "references", "citation", "citations", "bibliography"],
    }
    lowered_prompt = instructions.lower()
    lines = []
    used_categories = set()
    for category, terms in categories.items():
        if any(term in lowered_prompt for term in terms):
            used_categories.add(category)
            related_values = []
            for signal_name, values in signals.items():
                signal_lower = signal_name.lower()
                if category in signal_lower or any(term in signal_lower for term in terms):
                    related_values.extend(values)
            exact_hits = sorted({term for term in terms if term in content.lower()})
            evidence = sorted(set(related_values + exact_hits))[:20]
            if evidence:
                lines.append(f"- `{category}`: found evidence ({', '.join(evidence[:10])})")
            else:
                lines.append(f"- `{category}`: not found automatically; manual review may be needed")

    for item in _condition_matches(content, instructions):
        term = item["condition_term"]
        if any(term in categories.get(category, []) for category in used_categories):
            continue
        status = "found" if item["occurrences"] else "not found"
        lines.append(f"- `{term}`: {status} ({item['occurrences']} occurrence(s))")
        if item["sample_context"]:
            lines.append(f"  - Context: {item['sample_context'][:260]}")
    return lines
def extract_file_content(settings: Settings, path: Path, max_chars: int = REPORT_MAX_CHARS) -> dict[str, Any]:
    path = workspace.ensure_inside_root(settings, path)
    if workspace.is_secret_path(path):
        raise workspace.WorkspaceError("Refusing to analyze secret or .git files.")
    if not path.exists() or not path.is_file():
        raise workspace.WorkspaceError("Document analysis requires an existing file.")

    content, kind = _extract_generic_content(path, max_chars)
    suffix = path.suffix.lower()
    return {
        "path": workspace.relative(settings, path),
        "kind": kind,
        "suffix": suffix,
        "language_or_format": LANGUAGE_BY_SUFFIX.get(suffix, kind),
        "size": path.stat().st_size,
        "content": content,
        "signals": _generic_code_signals(content),
    }


def _line_block(title: str, values: list[str], empty: str) -> str:
    if values:
        return f"### {title}\n" + "\n".join(f"- {value}" for value in values[:40])
    return f"### {title}\n- {empty}"


def _section_text_for_kind(kind: str, suffix: str) -> dict[str, str]:
    if kind == "jupyter_notebook":
        return {
            "file_purpose": "The file is a Jupyter notebook and should be treated as an executable workflow, experiment, or analysis artifact. The extracted content includes markdown cells, code cells, available text outputs, and detected inline figure outputs when present.",
            "workflow": "The notebook workflow should be interpreted by grouping cells into meaningful stages such as data loading, inspection, preprocessing, model/system logic, evaluation, visualization, and reporting. Individual cells should not be treated as independent sections unless required by the evidence.",
        }
    if "source" in kind or kind in {"text_or_code", "unknown_text_like_file"}:
        return {
            "file_purpose": "The file appears to be a source-code or text-based technical file. Its purpose should be inferred from imports, functions, classes, commands, file operations, comments, and visible execution logic.",
            "workflow": "The execution flow should be interpreted from the visible code structure, including imports, configuration, inputs, main functions/classes, control flow, outputs, error handling, and integration points.",
        }
    if kind in {"pdf", "word_document", "powerpoint_presentation", "powerpoint_presentation_xml_fallback"}:
        return {
            "file_purpose": "The file appears to be a technical document or presentation. Its purpose should be inferred from headings, sections, tables, claims, methodology, results, and conclusions visible in extracted text.",
            "workflow": "The document flow should be described through its section structure, methodology, argument sequence, findings, conclusions, and any visible technical evidence.",
        }
    if kind == "spreadsheet":
        return {
            "file_purpose": "The file appears to be a spreadsheet. Its purpose should be inferred from sheet names, column headers, visible rows, tables, and any detected numeric or categorical data patterns.",
            "workflow": "The spreadsheet should be interpreted through sheets, columns, records, visible calculations, table layout, and possible data preparation or reporting role.",
        }
    return {
        "file_purpose": f"The file type `{suffix or '[no extension]'}` was analyzed using generic extraction. The purpose should be inferred only from available readable evidence.",
        "workflow": "The workflow or structure should be described from extracted readable content. Unsupported binary sections should not be overinterpreted.",
    }


def _line_block(title: str, values: list[str], empty: str) -> str:
    if values:
        return f"### {title}\n" + "\n".join(f"- {value}" for value in values[:40])
    return f"### {title}\n- {empty}"


def _section_text_for_kind(kind: str, suffix: str) -> dict[str, str]:
    if kind == "jupyter_notebook":
        return {
            "file_purpose": "The file is a Jupyter notebook and should be treated as an executable workflow, experiment, or analysis artifact. The extracted content includes markdown cells, code cells, available text outputs, and detected inline figure outputs when present.",
            "workflow": "The notebook workflow should be interpreted by grouping cells into meaningful stages such as data loading, inspection, preprocessing, model/system logic, evaluation, visualization, and reporting. Individual cells should not be treated as independent sections unless required by the evidence.",
        }
    if "source" in kind or kind in {"text_or_code", "unknown_text_like_file"}:
        return {
            "file_purpose": "The file appears to be a source-code or text-based technical file. Its purpose should be inferred from imports, functions, classes, commands, file operations, comments, and visible execution logic.",
            "workflow": "The execution flow should be interpreted from the visible code structure, including imports, configuration, inputs, main functions/classes, control flow, outputs, error handling, and integration points.",
        }
    if kind in {"pdf", "word_document", "powerpoint_presentation", "powerpoint_presentation_xml_fallback"}:
        return {
            "file_purpose": "The file appears to be a technical document or presentation. Its purpose should be inferred from headings, sections, tables, claims, methodology, results, and conclusions visible in extracted text.",
            "workflow": "The document flow should be described through its section structure, methodology, argument sequence, findings, conclusions, and any visible technical evidence.",
        }
    if kind == "spreadsheet":
        return {
            "file_purpose": "The file appears to be a spreadsheet. Its purpose should be inferred from sheet names, column headers, visible rows, tables, and any detected numeric or categorical data patterns.",
            "workflow": "The spreadsheet should be interpreted through sheets, columns, records, visible calculations, table layout, and possible data preparation or reporting role.",
        }
    return {
        "file_purpose": f"The file type `{suffix or '[no extension]'}` was analyzed using generic extraction. The purpose should be inferred only from available readable evidence.",
        "workflow": "The workflow or structure should be described from extracted readable content. Unsupported binary sections should not be overinterpreted.",
    }


def _wants_full_report(instructions: str) -> bool:
    lowered = instructions.lower()
    return any(term in lowered for term in ["full report", "detailed report", "technical report", "thesis", "project report", "complete report"])


def _append_section(sections: list[str], title: str, body: str, include: bool = True) -> None:
    if include and body.strip():
        sections.append(f"## {title}\n\n{body.strip()}")


def build_technical_report(settings: Settings, path: Path, instructions: str = "") -> str:
    extracted = extract_file_content(settings, path)
    rel = extracted["path"]
    content = extracted["content"]
    suffix = extracted["suffix"]
    kind = extracted["kind"]
    focus = instructions.strip() or "Create a technical summary from the analyzed file."
    signals = extracted.get("signals", {})
    condition_lines = _condition_lines_from_signals(content, focus, signals)
    kind_text = _section_text_for_kind(kind, suffix)
    full_report = _wants_full_report(focus)

    functions = signals.get("functions_or_methods", [])
    classes = signals.get("classes", [])
    dependencies = signals.get("imports_or_dependencies", [])
    preprocessing = signals.get("preprocessing_terms", [])
    models = signals.get("models_or_algorithms", [])
    metrics = signals.get("metrics_or_results", [])
    figures = signals.get("figures_or_plots", [])
    has_tables = "[Table" in content or "--- Sheet:" in content
    has_outputs = "### Output From Cell" in content or bool(metrics)

    sections = [
        "# Technical Summary Report",
        (
            "## Overview\n\n"
            f"The file contains visible technical content in `{extracted['language_or_format']}` format. "
            "The content is organized around the detected workflow, inputs, methods, outputs, and limitations available from the extracted evidence.\n\n"
            f"**Requested focus:** {focus}"
        ),
        (
            "## File Type and Purpose\n\n"
            f"- **Path:** `{rel}`\n"
            f"- **File type:** `{suffix or '[no extension]'}`\n"
            f"- **Detected format:** `{kind}`\n"
            f"- **Size:** {extracted['size']} bytes\n"
            f"- **Extracted text lines:** {content.count(chr(10)) + 1 if content else 0}\n\n"
            f"{kind_text['file_purpose']}"
        ),
    ]

    _append_section(
        sections,
        "Prompt Condition Checklist",
        "\n".join(condition_lines) if condition_lines else "No specific prompt condition keywords were extracted from the request.",
        include=bool(condition_lines),
    )

    _append_section(
        sections,
        "Inputs, Dependencies, and Setup",
        _line_block("Visible Dependencies Or Inputs", dependencies, "No explicit dependencies, imports, APIs, configuration files, datasets, or setup requirements were identified."),
        include=full_report or bool(dependencies) or kind in {"jupyter_notebook", "spreadsheet"},
    )

    _append_section(
        sections,
        "Workflow or Execution Summary",
        kind_text["workflow"],
        include=full_report or kind != "binary_or_unsupported_file",
    )

    preprocessing_body = "No explicit preprocessing or input preparation stage was identified."
    if preprocessing:
        preprocessing_body = "Detected preprocessing or input-handling terms include: " + ", ".join(preprocessing[:30]) + ". Only techniques supported by visible evidence should be claimed."
    _append_section(
        sections,
        "Data Preparation, Preprocessing, or Input Handling",
        preprocessing_body,
        include=full_report or bool(preprocessing) or any(term in focus.lower() for term in ["preprocess", "input", "data", "clean"]),
    )

    method_lines = []
    if functions:
        method_lines.append(_line_block("Functions Or Methods", functions, "No explicit functions or methods were detected."))
    if classes:
        method_lines.append(_line_block("Classes", classes, "No explicit classes were detected."))
    if models:
        method_lines.append(_line_block("Models, Algorithms, Or System Logic Terms", models, "No explicit model or algorithm names were detected."))
    _append_section(
        sections,
        "Methods, Models, Algorithms, or System Logic",
        "\n\n".join(method_lines) if method_lines else "Not explicitly available in the provided file.",
        include=full_report or bool(method_lines) or any(term in focus.lower() for term in ["model", "method", "algorithm", "function", "class", "logic"]),
    )

    outputs_statement = "Not explicitly available in the provided file."
    if metrics:
        outputs_statement = "Potential result or evaluation terms were detected: " + ", ".join(metrics[:30]) + ". Exact numeric values should be used only if shown in the extracted evidence."
    if "### Output From Cell" in content:
        outputs_statement = "Notebook text outputs are present in the extracted content. Reported values should be copied only from visible outputs."
    _append_section(
        sections,
        "Results, Outputs, and Evaluation",
        outputs_statement,
        include=full_report or has_outputs or any(term in focus.lower() for term in ["result", "metric", "evaluation", "output"]),
    )

    visual_statement = "No visual or tabular outputs were identified."
    if figures:
        visual_statement = "Potential visual output logic was identified through plotting or figure-related terms: " + ", ".join(figures[:20]) + "."
    if has_tables:
        visual_statement = "Tabular content was identified in the extracted file content."
    _append_section(
        sections,
        "Visual Outputs, Tables, and Figures",
        visual_statement,
        include=full_report or bool(figures) or has_tables or any(term in focus.lower() for term in ["figure", "plot", "table", "visual"]),
    )

    _append_section(
        sections,
        "Technical Interpretation",
        "The file should be interpreted only through visible evidence. Completeness, technical contribution, and result implications should be discussed only when supported by extracted code, outputs, document claims, or detected signals. If evidence is incomplete, conclusions cannot be fully verified from the available content.",
        include=full_report,
    )

    _append_section(
        sections,
        "Limitations, Risks, or Issues",
        "Potential limitations include missing datasets, hardcoded paths, missing dependencies, unexecuted notebook cells, missing outputs, incomplete evaluation, weak preprocessing evidence, lack of validation, missing error handling, ambiguous target variables, missing documentation, code failures, missing citations, or insufficient methodology details. Items not visible in the analyzed content should be reported as: ÃƒÆ’Ã†â€™Ãƒâ€šÃ‚Â¢ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â€šÂ¬Ã…Â¡Ãƒâ€šÃ‚Â¬ÃƒÆ’Ã¢â‚¬Â¦ÃƒÂ¢Ã¢â€šÂ¬Ã…â€œNot explicitly available in the provided file.ÃƒÆ’Ã†â€™Ãƒâ€šÃ‚Â¢ÃƒÆ’Ã‚Â¢ÃƒÂ¢Ã¢â€šÂ¬Ã…Â¡Ãƒâ€šÃ‚Â¬ÃƒÆ’Ã¢â‚¬Å¡Ãƒâ€šÃ‚Â",
        include=full_report or any(term in focus.lower() for term in ["limitation", "risk", "issue", "missing"]),
    )

    _append_section(
        sections,
        "Recommended Improvements",
        "Recommended improvements should be practical and evidence-based: add documentation where workflow steps are unclear; document dependencies; save generated outputs or metrics; add validation, testing, baselines, configuration files, reproducibility instructions, and error handling where relevant.",
        include=full_report or any(term in focus.lower() for term in ["recommend", "improve", "suggest"]),
    )

    _append_section(
        sections,
        "Final Summary",
        f"The analyzed file is `{rel}`, detected as `{extracted['language_or_format']}`. Its technical value should be summarized in terms of visible purpose, workflow, methods, outputs or results, limitations, and evidence-backed conclusions. Missing information should be clearly identified rather than inferred.",
        include=full_report,
    )

    _append_section(
        sections,
        "Appendix: Extracted Evidence",
        f"```text\n{content[:28000]}\n```",
        include=full_report or len(content) < 12000,
    )

    return "\n\n".join(sections)

# DOCX style-preserving editing helpers

def _style_names(doc) -> set[str]:
    return {style.name for style in doc.styles}


def _first_available_style(doc, candidates: list[str], fallback: str = "Normal") -> str:
    names = _style_names(doc)
    for candidate in candidates:
        if candidate in names:
            return candidate
    return fallback if fallback in names else "Normal"


def _body_style(doc) -> str:
    counts = Counter()
    for paragraph in doc.paragraphs:
        if not paragraph.text.strip() or not paragraph.style:
            continue
        name = paragraph.style.name
        lowered = name.lower()
        if "heading" in lowered or "title" in lowered or "list" in lowered:
            continue
        counts[name] += 1
    if counts:
        return counts.most_common(1)[0][0]
    return _first_available_style(doc, ["Body Text", "Normal"], "Normal")


def _style_map(doc, style_hint: str = "") -> dict[str, str]:
    names = _style_names(doc)
    requested = style_hint.strip()
    body = requested if requested in names else _body_style(doc)
    return {
        "body": body,
        "heading_1": _first_available_style(doc, ["Heading 1", "Title"], body),
        "heading_2": _first_available_style(doc, ["Heading 2", "Subtitle", "Heading 1"], body),
        "heading_3": _first_available_style(doc, ["Heading 3", "Heading 2", "Heading 1"], body),
        "bullet": _first_available_style(doc, ["List Bullet", "List Paragraph"], body),
        "number": _first_available_style(doc, ["List Number", "List Paragraph"], body),
    }


def _apply_style(paragraph, style_name: str) -> None:
    try:
        paragraph.style = style_name
    except (KeyError, ValueError):
        pass


def _strip_inline_markdown(text: str) -> str:
    text = _sanitize_xml_text(text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"__(.*?)__", r"\1", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    return text.strip()


def _paragraph_kind_from_style(paragraph) -> str:
    style_name = paragraph.style.name.lower() if paragraph.style else ""
    if "heading 1" in style_name or style_name == "title":
        return "heading_1"
    if "heading 2" in style_name or "subtitle" in style_name:
        return "heading_2"
    if "heading 3" in style_name:
        return "heading_3"
    return "body"


def _add_markdown_line(doc, line: str, styles: dict[str, str]) -> str | None:
    stripped = line.strip()
    if not stripped:
        return None
    heading = re.match(r"^(#{1,6})\s+(.+)$", stripped)
    if heading:
        level = min(len(heading.group(1)), 6)
        if level == 1 and not styles.get("_title_used"):
            kind = "title"
            style_key = "heading_1"
            styles["_title_used"] = True
        elif level <= 2:
            kind = "heading_1"
            style_key = "heading_1"
        else:
            kind = "heading_2"
            style_key = "heading_2"
        paragraph = doc.add_paragraph()
        _apply_style(paragraph, styles[style_key])
        paragraph.add_run(_strip_inline_markdown(heading.group(2)))
        _format_paragraph(paragraph, kind)
        return kind
    bullet = re.match(r"^[-*+]\s+(.+)$", stripped)
    if bullet:
        paragraph = doc.add_paragraph()
        _apply_style(paragraph, styles["bullet"])
        paragraph.add_run(_strip_inline_markdown(bullet.group(1)))
        _format_paragraph(paragraph, "body")
        return "bullet"
    numbered = re.match(r"^\d+[.)]\s+(.+)$", stripped)
    if numbered:
        paragraph = doc.add_paragraph()
        _apply_style(paragraph, styles["number"])
        paragraph.add_run(_strip_inline_markdown(numbered.group(1)))
        _format_paragraph(paragraph, "body")
        return "number"
    paragraph = doc.add_paragraph()
    _apply_style(paragraph, styles["body"])
    paragraph.add_run(_strip_inline_markdown(stripped))
    _format_paragraph(paragraph, "body")
    return "body"

def write_docx_content(settings: Settings, path: Path, content: str, overwrite: bool = False, style_hint: str = "") -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentToolError("DOCX writing requires python-docx. Install dependencies from requirements.txt.") from exc

    path = workspace.ensure_inside_root(settings, path)
    if path.suffix.lower() != ".docx":
        raise workspace.WorkspaceError("Word report generation requires a .docx path.")
    if workspace.is_secret_path(path):
        raise workspace.WorkspaceError("Refusing to write secret or .git files.")
    if path.exists() and not overwrite:
        raise workspace.WorkspaceError("DOCX file already exists. Set overwrite=True or choose a new path.")

    backup = None
    if path.exists():
        backup = workspace.make_backup(settings, path)

    path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    styles = _style_map(doc, style_hint=style_hint)
    inserted = Counter()
    for line in content.splitlines():
        kind = _add_markdown_line(doc, line, styles)
        if kind:
            inserted[kind] += 1
    doc.save(str(path))
    return {
        "path": workspace.relative(settings, path),
        "backup": workspace.relative(settings, backup) if backup else None,
        "mode": "write_docx_markdown_preserving_defaults",
        "styles_used": {key: value for key, value in styles.items() if not key.startswith("_")},
        "inserted": dict(inserted),
    }

def append_docx_content(settings: Settings, path: Path, content: str, style_hint: str = "") -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentToolError("DOCX editing requires python-docx. Install dependencies from requirements.txt.") from exc

    path = workspace.ensure_inside_root(settings, path)
    if path.suffix.lower() != ".docx":
        raise workspace.WorkspaceError("Style-preserving document editing currently supports .docx files.")
    if workspace.is_secret_path(path):
        raise workspace.WorkspaceError("Refusing to edit secret or .git files.")

    backup = workspace.make_backup(settings, path)
    doc = Document(str(path))
    styles = _style_map(doc, style_hint=style_hint)
    inserted = Counter()
    for line in content.splitlines():
        kind = _add_markdown_line(doc, line, styles)
        if kind:
            inserted[kind] += 1
    doc.save(str(path))
    return {
        "path": workspace.relative(settings, path),
        "backup": workspace.relative(settings, backup),
        "mode": "append_markdown_preserving_styles",
        "styles_used": {key: value for key, value in styles.items() if not key.startswith("_")},
        "inserted": dict(inserted),
    }


def replace_docx_paragraph(settings: Settings, path: Path, marker: str, replacement: str) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentToolError("DOCX editing requires python-docx. Install dependencies from requirements.txt.") from exc

    path = workspace.ensure_inside_root(settings, path)
    if path.suffix.lower() != ".docx":
        raise workspace.WorkspaceError("Style-preserving document editing currently supports .docx files.")
    if not marker:
        raise workspace.WorkspaceError("Replacing DOCX content requires a marker or existing paragraph text to find.")

    backup = workspace.make_backup(settings, path)
    doc = Document(str(path))
    changed = False
    for paragraph in doc.paragraphs:
        if marker in paragraph.text:
            style = paragraph.style
            run_style = paragraph.runs[0] if paragraph.runs else None
            paragraph.clear()
            paragraph.style = style
            run = paragraph.add_run(_strip_inline_markdown(replacement))
            _format_paragraph(paragraph, "body")
            if run_style is not None:
                run.bold = run_style.bold
                run.italic = run_style.italic
                run.underline = run_style.underline
                if run_style.font is not None:
                    run.font.name = run_style.font.name
                    run.font.size = run_style.font.size
            changed = True
            break
    if not changed:
        raise workspace.WorkspaceError("Could not find the requested marker text in the DOCX file.")
    doc.save(str(path))
    return {"path": workspace.relative(settings, path), "backup": workspace.relative(settings, backup), "mode": "replace", "marker": marker}


def modify_docx(settings: Settings, path: Path, mode: str, content: str, marker: str = "", style_hint: str = "") -> dict[str, Any]:
    if mode in {"append", "append_markdown"}:
        return append_docx_content(settings, path, content, style_hint=style_hint)
    if mode == "replace_paragraph":
        return replace_docx_paragraph(settings, path, marker, content)
    raise workspace.WorkspaceError("Unsupported DOCX modification mode. Use append, append_markdown, or replace_paragraph.")













# Notebook-specific technical report generation

_PREPROCESSING_PATTERNS = [
    (r"\bfillna\s*\(|\bSimpleImputer\b", "Missing-value imputation"),
    (r"\bdropna\s*\(", "Missing-value row/column removal"),
    (r"\binterpolate\s*\(", "Interpolation"),
    (r"\bStandardScaler\b|\bstandardscaler\s*\(", "Standard scaling"),
    (r"\bMinMaxScaler\b|\bminmaxscaler\s*\(", "Min-max scaling"),
    (r"\bRobustScaler\b", "Robust scaling"),
    (r"\bLabelEncoder\b", "Label encoding"),
    (r"\bOneHotEncoder\b|\bget_dummies\s*\(", "One-hot encoding"),
    (r"\btrain_test_split\s*\(", "Train/test split"),
    (r"\bto_datetime\s*\(", "Datetime conversion"),
    (r"\bset_index\s*\(", "Index assignment"),
    (r"\bresample\s*\(", "Time-series resampling"),
    (r"\brolling\s*\(", "Rolling-window feature calculation"),
    (r"\bshift\s*\(", "Lag/shift feature engineering"),
    (r"\breshape\s*\(|\bTimeseriesGenerator\b|\bsequence\b|\blook_back\b", "Time-series sequence/window preparation"),
    (r"\bdrop_duplicates\s*\(", "Duplicate removal"),
    (r"\bastype\s*\(", "Data type conversion"),
    (r"\bconcat\s*\(|\bmerge\s*\(|\bjoin\s*\(", "Feature/table merging"),
]

_MODEL_PATTERNS = [
    (r"\bLinearRegression\b", "Linear Regression"),
    (r"\bLogisticRegression\b", "Logistic Regression"),
    (r"\bLSTM\b", "LSTM"),
    (r"\bGRU\b", "GRU"),
    (r"\bSequential\b", "Keras Sequential neural network"),
    (r"\bDense\b", "Dense neural network layers"),
    (r"\bRandomForestRegressor\b", "Random Forest Regressor"),
    (r"\bRandomForestClassifier\b", "Random Forest Classifier"),
    (r"\bDecisionTreeRegressor\b", "Decision Tree Regressor"),
    (r"\bDecisionTreeClassifier\b", "Decision Tree Classifier"),
    (r"\bSVR\b", "Support Vector Regressor"),
    (r"\bSVC\b", "Support Vector Classifier"),
    (r"\bXGBRegressor\b|\bXGBoost\b", "XGBoost"),
    (r"\bARIMA\b|\bSARIMA\b", "ARIMA/SARIMA"),
    (r"\bProphet\b", "Prophet"),
    (r"\bKNeighbors\w*\b", "K-Nearest Neighbors"),
]

_METRIC_ALIASES = {
    "mae": "MAE",
    "mean absolute error": "MAE",
    "mse": "MSE",
    "mean squared error": "MSE",
    "rmse": "RMSE",
    "root mean squared error": "RMSE",
    "r2": "R2",
    "r^2": "R2",
    "r-squared": "R2",
    "accuracy": "Accuracy",
    "precision": "Precision",
    "recall": "Recall",
    "f1": "F1 Score",
    "loss": "Loss",
    "val_loss": "Validation Loss",
    "score": "Score",
}


def _nb_source(cell: dict[str, Any]) -> str:
    source = cell.get("source", "")
    if isinstance(source, list):
        return "".join(str(item) for item in source)
    return str(source or "")


def _nb_output_text(cell: dict[str, Any]) -> str:
    parts = []
    for output in cell.get("outputs", []):
        if output.get("ename") or output.get("evalue"):
            parts.append(f"Error: {output.get('ename', '')} {output.get('evalue', '')}".strip())
        text_items = output.get("text") or output.get("data", {}).get("text/plain") or []
        if isinstance(text_items, str):
            text_items = [text_items]
        if text_items:
            parts.append(_sanitize_xml_text("".join(str(item) for item in text_items)).strip())
    return "\n".join(item for item in parts if item)


def _detect_patterns(source: str, patterns: list[tuple[str, str]]) -> list[str]:
    found = []
    for pattern, label in patterns:
        if re.search(pattern, source, flags=re.IGNORECASE):
            found.append(label)
    return found


def _extract_imports(source: str) -> list[str]:
    imports = []
    for match in re.findall(r"^(?:import|from)\s+([^\n#]+)", source, flags=re.MULTILINE):
        item = match.strip()
        if item and item not in imports:
            imports.append(item[:120])
    return imports


def _extract_inputs(source: str) -> list[str]:
    inputs = []
    for pattern in [r"read_csv\(([^\)]*)\)", r"read_excel\(([^\)]*)\)", r"read_table\(([^\)]*)\)", r"load_workbook\(([^\)]*)\)"]:
        for match in re.findall(pattern, source, flags=re.IGNORECASE):
            cleaned = _strip_inline_markdown(match).strip().strip(",")
            if cleaned and cleaned not in inputs:
                inputs.append(cleaned[:160])
    return inputs


def _extract_metrics(output_text: str) -> list[dict[str, str]]:
    metrics = []
    if not output_text:
        return metrics
    name_pattern = r"(?:mean absolute error|mean squared error|root mean squared error|r-squared|val_loss|accuracy|precision|recall|rmse|mae|mse|r\^2|r2|f1|loss|score)"
    value_pattern = r"[-+]?\d*\.?\d+(?:[eE][-+]?\d+)?%?"
    for line in output_text.splitlines():
        line_clean = re.sub(r"\s+", " ", line).strip()
        if not line_clean:
            continue
        for match in re.finditer(rf"({name_pattern})\s*[:=\-]?\s*({value_pattern})", line_clean, flags=re.IGNORECASE):
            raw_name = match.group(1).lower()
            metric = _METRIC_ALIASES.get(raw_name, raw_name.upper())
            value = match.group(2)
            entry = {"metric": metric, "value": value, "line": line_clean[:220]}
            if entry not in metrics:
                metrics.append(entry)
    return metrics


def _extract_notebook_figures(settings: Settings, path: Path) -> list[dict[str, str]]:
    import base64

    raw = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    out_dir = workspace.ensure_inside_root(settings, path.parent / f"{path.stem}_figures")
    figures = []
    for cell_index, cell in enumerate(raw.get("cells", []), start=1):
        for output_index, output in enumerate(cell.get("outputs", []), start=1):
            data = output.get("data", {})
            image_key = None
            extension = None
            for candidate, ext in (("image/png", "png"), ("image/jpeg", "jpg"), ("image/jpg", "jpg")):
                if candidate in data:
                    image_key = candidate
                    extension = ext
                    break
            if not image_key:
                continue
            encoded = data.get(image_key)
            if isinstance(encoded, list):
                encoded = "".join(encoded)
            if not encoded:
                continue
            out_dir.mkdir(parents=True, exist_ok=True)
            filename = f"figure_cell_{cell_index:03d}_output_{output_index:02d}.{extension}"
            target = workspace.ensure_inside_root(settings, out_dir / filename)
            target.write_bytes(base64.b64decode(encoded))
            figures.append({
                "cell": str(cell_index),
                "path": workspace.relative(settings, target),
                "absolute_path": str(target),
            })
    return figures


def _cell_role(source: str, preprocessing: list[str], models: list[str], imports: list[str], inputs: list[str], metrics: list[dict[str, str]], has_figures: bool) -> str:
    lowered = source.lower()
    roles = []
    if imports:
        roles.append("imports dependencies")
    if inputs or "read_csv" in lowered or "read_excel" in lowered:
        roles.append("loads data")
    if any(term in lowered for term in ["head(", "info(", "describe(", "isnull", "shape", "columns"]):
        roles.append("inspects the dataset")
    if preprocessing:
        roles.append("applies preprocessing")
    if models:
        roles.append("defines or trains model logic")
    if any(term in lowered for term in ["predict", "evaluate", "score", "mean_squared_error", "mean_absolute_error", "r2_score"]):
        roles.append("evaluates predictions")
    if has_figures or any(term in lowered for term in ["plot", "figure", "imshow", "scatter", "hist", "bar"]):
        roles.append("creates visual output")
    if metrics:
        roles.append("prints numerical results")
    return ", ".join(roles) if roles else "contains code with no automatically classified workflow role"


def _group_by_name(items: list[dict[str, Any]], key: str) -> list[str]:
    grouped: dict[str, set[str]] = {}
    for item in items:
        grouped.setdefault(item[key], set()).add(str(item["cell"]))
    return [f"- {name}: detected in code cell(s) {', '.join(sorted(cells, key=lambda value: int(value)))}" for name, cells in sorted(grouped.items())]


def _build_notebook_technical_report(settings: Settings, path: Path, instructions: str = "", include_figures: bool = True) -> str:
    raw = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    cells = raw.get("cells", [])
    rel = workspace.relative(settings, path)
    focus = instructions.strip() or "Create a technical report from the notebook."

    embedded_figure_count = sum(1 for cell in cells for output in cell.get("outputs", []) if "data" in output and any(str(key).startswith("image/") for key in output.get("data", {})))
    figures = _extract_notebook_figures(settings, path) if include_figures else []
    figure_cells = {item["cell"] for item in figures}
    code_summaries = []
    preprocessing_items = []
    model_items = []
    metric_items = []
    imports = []
    inputs = []
    executed = 0
    cells_with_outputs = 0
    errors = []

    for index, cell in enumerate(cells, start=1):
        if cell.get("cell_type") != "code":
            continue
        source = _nb_source(cell)
        output_text = _nb_output_text(cell)
        cell_imports = _extract_imports(source)
        cell_inputs = _extract_inputs(source)
        cell_preprocessing = _detect_patterns(source, _PREPROCESSING_PATTERNS)
        cell_models = _detect_patterns(source, _MODEL_PATTERNS)
        cell_metrics = _extract_metrics(output_text)
        if cell.get("execution_count") is not None:
            executed += 1
        if cell.get("outputs"):
            cells_with_outputs += 1
        for item in cell_imports:
            if item not in imports:
                imports.append(item)
        for item in cell_inputs:
            if item not in inputs:
                inputs.append(item)
        for item in cell_preprocessing:
            preprocessing_items.append({"name": item, "cell": index})
        for item in cell_models:
            model_items.append({"name": item, "cell": index})
        for metric in cell_metrics:
            metric_items.append({**metric, "cell": index})
        if "Error:" in output_text:
            errors.append(f"- Cell {index}: {output_text.splitlines()[0][:180]}")
        code_summaries.append(
            {
                "cell": index,
                "role": _cell_role(source, cell_preprocessing, cell_models, cell_imports, cell_inputs, cell_metrics, str(index) in figure_cells),
                "preprocessing": cell_preprocessing,
                "models": cell_models,
                "metrics": cell_metrics,
                "has_figure": str(index) in figure_cells,
            }
        )

    unique_preprocessing = sorted({item["name"] for item in preprocessing_items})
    unique_models = sorted({item["name"] for item in model_items})
    unique_metric_names = sorted({item["metric"] for item in metric_items})
    metric_labels = sorted({item.get("label", "") for item in metric_items if item.get("label")})

    input_text = ", ".join(inputs[:3]) if inputs else "the dataset loaded inside the notebook"
    preprocessing_text = ", ".join(unique_preprocessing[:8]) if unique_preprocessing else "no clearly detected preprocessing step"
    model_text = ", ".join(unique_models[:6]) if unique_models else "no clearly detected model"
    metric_text = ", ".join(unique_metric_names[:6]) if unique_metric_names else "no parseable metric names"
    experiment_text = ", ".join(metric_labels[:8]) if metric_labels else "the visible notebook experiments"
    figure_text = f"The notebook also contains {embedded_figure_count} embedded visual output(s)." if embedded_figure_count else "No embedded visual outputs were detected."
    overview_paragraph = (
        f"The notebook loads {input_text}, inspects the dataset, and runs an experimental machine-learning workflow. "
        f"The workflow applies preprocessing and feature-preparation steps including {preprocessing_text}. "
        f"It then trains or evaluates {model_text}, with evaluation output reported using {metric_text}. "
        f"The visible result labels indicate experiments such as {experiment_text}. "
        f"{figure_text} The notebook is therefore organized around dataset preparation, model evaluation, result comparison, and visualization."
    )

    sections = [
        "# Technical Summary Report",
        "## Overview\n\n" + overview_paragraph,
        (
            "## Notebook Execution Profile\n\n"
            f"- Total cells: {len(cells)}\n"
            f"- Code cells: {sum(1 for cell in cells if cell.get('cell_type') == 'code')}\n"
            f"- Markdown cells: {sum(1 for cell in cells if cell.get('cell_type') == 'markdown')}\n"
            f"- Executed code cells: {executed}\n"
            f"- Code cells with outputs: {cells_with_outputs}\n"
            f"- Embedded figures detected: {embedded_figure_count}\n"
            f"- Figure files extracted: {len(figures)}\n"
            f"- Requested focus: {focus}"
        ),
    ]

    if imports or inputs:
        body = []
        if imports:
            body.append("### Libraries and Dependencies\n" + "\n".join(f"- {item}" for item in imports[:40]))
        if inputs:
            body.append("### Visible Data Inputs\n" + "\n".join(f"- {item}" for item in inputs[:20]))
        sections.append("## Inputs, Dependencies, and Setup\n\n" + "\n\n".join(body))

    workflow_lines = []
    for item in code_summaries:
        details = [f"Cell {item['cell']}: {item['role']}."]
        if item["preprocessing"]:
            details.append("Preprocessing: " + ", ".join(item["preprocessing"]) + ".")
        if item["models"]:
            details.append("Models/system logic: " + ", ".join(item["models"]) + ".")
        if item["metrics"]:
            details.append("Visible metrics: " + ", ".join(f"{m['metric']}={m['value']}" for m in item["metrics"][:8]) + ".")
        if item["has_figure"]:
            details.append("Embedded figure output detected.")
        workflow_lines.append("- " + " ".join(details))
    sections.append("## Workflow Summary by Code Cell\n\n" + ("\n".join(workflow_lines) if workflow_lines else "No code cells were available for analysis."))

    sections.append(
        "## Data Preparation and Preprocessing Techniques\n\n"
        + ("\n".join(_group_by_name(preprocessing_items, "name")) if preprocessing_items else "No explicit preprocessing technique was detected in the code cells.")
    )

    sections.append(
        "## Models and Algorithms Used\n\n"
        + ("\n".join(_group_by_name(model_items, "name")) if model_items else "No explicit model or algorithm name was detected in the code cells.")
    )

    if metric_items:
        result_lines = []
        for metric in metric_items:
            result_lines.append(f"- Cell {metric['cell']}: {metric['metric']} = {metric['value']} ({metric['line']})")
        results_body = "\n".join(result_lines)
    else:
        results_body = "No numerical metric values were extracted from visible cell text outputs. If metrics are present only inside images, they should be read from the saved figures manually."
    sections.append("## Results and Evaluation Metrics\n\n" + results_body)

    if figures:
        figure_lines = []
        for number, figure in enumerate(figures, start=1):
            figure_lines.append(f"- Figure {number}: output from code cell {figure['cell']} saved at `{figure['path']}`")
            figure_lines.append(f"![Figure {number}]({figure['absolute_path']})")
        sections.append("## Visual Outputs and Figures\n\n" + "\n".join(figure_lines))
    else:
        sections.append("## Visual Outputs and Figures\n\nNo embedded notebook figures were detected in the analyzed cell outputs.")

    interpretation = []
    if preprocessing_items:
        interpretation.append("The notebook includes explicit preprocessing steps before or around modeling.")
    if model_items:
        interpretation.append("The notebook applies identifiable modeling logic: " + ", ".join(sorted({item['name'] for item in model_items})) + ".")
    if metric_items:
        interpretation.append("Visible output metrics provide direct evaluation evidence and are listed in the results section.")
    if not interpretation:
        interpretation.append("The available notebook evidence is limited, so technical conclusions should remain conservative.")
    sections.append("## Technical Interpretation\n\n" + " ".join(interpretation))

    limitation_lines = []
    if executed < sum(1 for cell in cells if cell.get("cell_type") == "code"):
        limitation_lines.append("- Some code cells appear unexecuted or have no execution count.")
    if not metric_items:
        limitation_lines.append("- No parseable numerical metrics were found in text outputs.")
    if errors:
        limitation_lines.extend(errors[:10])
    if not limitation_lines:
        limitation_lines.append("- No execution errors were detected from visible text outputs. Conclusions still depend on the completeness of the notebook outputs.")
    sections.append("## Limitations and Issues\n\n" + "\n".join(limitation_lines))

    sections.append(
        "## Recommended Improvements\n\n"
        "- Keep final metrics in clearly labeled output cells or exported result tables.\n"
        "- Save important plots with descriptive filenames.\n"
        "- Add markdown explanations for dataset columns, target variable, model configuration, and evaluation protocol.\n"
        "- Add reproducibility details, including dependencies, random seeds, and data source paths."
    )

    sections.append(
        "## Final Summary\n\n"
        f"The notebook `{rel}` was analyzed cell by cell. The generated report identifies detected preprocessing techniques, model names, visible evaluation metrics, and extracted figures without copying code cells into the report body."
    )
    return "\n\n".join(sections)


_previous_build_technical_report = build_technical_report


def build_technical_report(settings: Settings, path: Path, instructions: str = "") -> str:
    path = workspace.ensure_inside_root(settings, path)
    if path.suffix.lower() == ".ipynb":
        return _build_notebook_technical_report(settings, path, instructions)
    return _previous_build_technical_report(settings, path, instructions)


_previous_add_markdown_line = _add_markdown_line


def _add_markdown_line(doc, line: str, styles: dict[str, str]) -> str | None:
    stripped = line.strip()
    image = re.match(r"^!\[(.*?)\]\((.*?)\)$", stripped)
    if image:
        caption = _sanitize_xml_text(image.group(1) or "Figure")
        image_path = Path(image.group(2).strip().strip('"'))
        paragraph = doc.add_paragraph()
        _apply_style(paragraph, styles["body"])
        paragraph.add_run(caption)
        _format_paragraph(paragraph, "body")
        if image_path.exists():
            try:
                from docx.shared import Inches

                doc.add_picture(str(image_path), width=Inches(6.0))
                return "figure"
            except Exception:
                note = doc.add_paragraph()
                _apply_style(note, styles["body"])
                note.add_run(f"Figure file saved at: {image_path}")
                _format_paragraph(note, "body")
                return "figure_reference"
        note = doc.add_paragraph()
        _apply_style(note, styles["body"])
        note.add_run(f"Figure file was not found at: {image_path}")
        _format_paragraph(note, "body")
        return "figure_missing"
    return _previous_add_markdown_line(doc, line, styles)

# Generic non-notebook report builders. These avoid raw source/document dumps.

_SOURCE_CODE_KINDS = {"text_or_code", "unknown_text_like_file"}
_DOCUMENT_KINDS = {"pdf", "word_document", "powerpoint_presentation", "powerpoint_presentation_xml_fallback"}


def _extract_code_structure(content: str, suffix: str) -> dict[str, list[str]]:
    structure: dict[str, list[str]] = {}
    signals = _generic_code_signals(content)
    structure.update(signals)
    if suffix in {".html", ".htm"}:
        tags = sorted(set(re.findall(r"<\s*([a-zA-Z][a-zA-Z0-9-]*)\b", content)))[:40]
        forms = re.findall(r"<\s*(form|input|button|select|textarea)\b", content, flags=re.IGNORECASE)
        scripts = re.findall(r"<\s*script\b[^>]*src=[\"']([^\"']+)[\"']", content, flags=re.IGNORECASE)
        if tags:
            structure["html_tags"] = tags
        if forms:
            structure["interactive_elements"] = sorted(set(item.lower() for item in forms))
        if scripts:
            structure["external_scripts"] = scripts[:20]
    if suffix in {".css", ".scss"}:
        selectors = re.findall(r"(^|\})\s*([^@{}][^{}]{0,160})\s*\{", content, flags=re.MULTILINE)
        colors = re.findall(r"#[0-9a-fA-F]{3,8}|rgba?\([^\)]+\)|hsla?\([^\)]+\)", content)
        media = re.findall(r"@media\s+([^\{]+)", content)
        if selectors:
            structure["css_selectors"] = [item[1].strip() for item in selectors[:40] if item[1].strip()]
        if colors:
            structure["css_colors"] = sorted(set(colors))[:30]
        if media:
            structure["responsive_rules"] = [item.strip() for item in media[:20]]
    if suffix in {".cpp", ".cc", ".cxx", ".c", ".h", ".hpp"}:
        includes = re.findall(r"#include\s+[<\"]([^>\"]+)[>\"]", content)
        namespaces = re.findall(r"using\s+namespace\s+([A-Za-z_][A-Za-z0-9_:]*)", content)
        if includes:
            structure["c_cpp_includes"] = sorted(set(includes))[:40]
        if namespaces:
            structure["namespaces"] = sorted(set(namespaces))[:20]
    return structure


def _list_or_missing(values: list[str], missing: str) -> str:
    cleaned = []
    for value in values:
        value = _sanitize_xml_text(value).strip()
        if value and value not in cleaned:
            cleaned.append(value)
    return "\n".join(f"- {value}" for value in cleaned[:50]) if cleaned else missing


def _build_source_code_report(settings: Settings, path: Path, instructions: str = "") -> str:
    extracted = extract_file_content(settings, path)
    rel = extracted["path"]
    suffix = extracted["suffix"]
    content = extracted["content"]
    focus = instructions.strip() or "Create a technical report from the source file."
    structure = _extract_code_structure(content, suffix)
    dependencies = structure.get("imports_or_dependencies", []) + structure.get("c_cpp_includes", [])
    functions = structure.get("functions_or_methods", [])
    classes = structure.get("classes", [])
    models = structure.get("models_or_algorithms", [])
    preprocessing = structure.get("preprocessing_terms", [])
    metrics = structure.get("metrics_or_results", [])
    figures = structure.get("figures_or_plots", [])

    purpose_parts = []
    if functions or classes:
        purpose_parts.append("defined functions/classes")
    if dependencies:
        purpose_parts.append("external dependencies")
    if models:
        purpose_parts.append("model or algorithm logic")
    if preprocessing:
        purpose_parts.append("data/input preparation logic")
    if suffix in {".html", ".htm"}:
        purpose_parts.append("web page structure")
    if suffix in {".css", ".scss"}:
        purpose_parts.append("visual styling rules")
    purpose = ", ".join(purpose_parts) if purpose_parts else "general source-code behavior"

    sections = [
        "# Technical Summary Report",
        (
            "## Overview\n\n"
            f"The file contains {purpose}. "
            f"It uses {', '.join(dependencies[:6]) if dependencies else 'no explicitly detected external dependency'} "
            "and includes visible structures such as detected functions, classes, markup, selectors, inputs, outputs, or integration points where present."
        ),
        (
            "## File Type and Purpose\n\n"
            f"- Path: `{rel}`\n"
            f"- File type: `{suffix or '[no extension]'}`\n"
            f"- Size: {extracted['size']} bytes\n"
            f"- Requested focus: {focus}"
        ),
        "## Dependencies and Inputs\n\n" + _list_or_missing(dependencies, "No explicit imports, includes, packages, APIs, or external dependencies were detected."),
    ]

    if suffix in {".html", ".htm"}:
        sections.append("## HTML Structure\n\n" + _list_or_missing(structure.get("html_tags", []), "No HTML tags were detected.") )
        sections.append("## Interactive or Linked Elements\n\n" + _list_or_missing(structure.get("interactive_elements", []) + structure.get("external_scripts", []), "No forms, controls, or external scripts were detected."))
    elif suffix in {".css", ".scss"}:
        sections.append("## Styling Structure\n\n" + _list_or_missing(structure.get("css_selectors", []), "No CSS selectors were detected."))
        sections.append("## Visual Design Signals\n\n" + _list_or_missing(structure.get("css_colors", []) + structure.get("responsive_rules", []), "No colors or responsive rules were detected."))
    else:
        sections.append("## Functions, Classes, and Main Logic\n\n" + _list_or_missing(functions + classes, "No explicit functions or classes were detected."))

    if preprocessing:
        sections.append("## Data Preparation or Input Handling\n\n" + _list_or_missing(preprocessing, "No preprocessing terms were detected."))
    if models:
        sections.append("## Models, Algorithms, or System Logic\n\n" + _list_or_missing(models, "No model or algorithm terms were detected."))
    if metrics or figures:
        sections.append("## Outputs, Metrics, and Visualizations\n\n" + _list_or_missing(metrics + figures, "No output, metric, or visualization signals were detected."))

    limitations = []
    if not dependencies:
        limitations.append("- Dependency requirements are not explicit in this file.")
    if "try" not in content.lower() and "except" not in content.lower() and suffix in CODE_AND_TEXT_SUFFIXES:
        limitations.append("- No obvious exception-handling block was detected.")
    if not functions and not classes and suffix not in {".html", ".htm", ".css", ".scss"}:
        limitations.append("- The file may be script-like or declarative, because no function/class structure was detected.")
    sections.append("## Limitations and Risks\n\n" + ("\n".join(limitations) if limitations else "No major structural limitation was automatically detected from the source text."))
    sections.append(
        "## Recommended Improvements\n\n"
        "- Add or improve comments/docstrings around important logic.\n"
        "- Document required inputs, outputs, and runtime assumptions.\n"
        "- Add error handling and tests where behavior is important.\n"
        "- Keep configuration and secrets outside source files."
    )
    sections.append(f"## Final Summary\n\nThe file `{rel}` was analyzed structurally. The report identifies detected dependencies, code structure, logic signals, outputs, risks, and improvements without reproducing the source code.")
    return "\n\n".join(sections)


def _document_section_candidates(content: str) -> list[str]:
    candidates = []
    for line in content.splitlines():
        clean = re.sub(r"\s+", " ", line).strip()
        if not clean or len(clean) > 140:
            continue
        if clean.startswith("--- ") or clean.startswith("[") or clean.istitle() or clean.isupper():
            if clean not in candidates:
                candidates.append(clean)
    return candidates[:40]


def _build_document_report(settings: Settings, path: Path, instructions: str = "") -> str:
    extracted = extract_file_content(settings, path)
    rel = extracted["path"]
    content = extracted["content"]
    focus = instructions.strip() or "Create a technical report from the document."
    signals = _generic_code_signals(content)
    sections_found = _document_section_candidates(content)
    condition_lines = _condition_lines_from_signals(content, focus, signals)
    tables = [line for line in content.splitlines() if " | " in line or line.startswith("[Table")][:30]
    methodology_terms = sorted(set(re.findall(r"\b(?:methodology|method|framework|model|algorithm|experiment|evaluation|result|finding|conclusion|dataset|analysis|architecture)\b", content, flags=re.IGNORECASE)))[:30]

    sections = [
        "# Technical Summary Report",
        (
            "## Overview\n\n"
            f"The document is structured around {', '.join(sections_found[:5]) if sections_found else 'the visible extracted text'}. "
            f"Detected technical signals include {', '.join(methodology_terms[:8]) if methodology_terms else 'no explicit methodology or result terms'}. "
            "The content indicates the document's topic, section flow, methodology evidence, tables, results, and limitations where those elements are visible in the extracted text."
        ),
        (
            "## File Type and Purpose\n\n"
            f"- Path: `{rel}`\n"
            f"- File type: `{extracted['suffix'] or '[no extension]'}`\n"
            f"- Size: {extracted['size']} bytes\n"
            f"- Requested focus: {focus}"
        ),
        "## Detected Document Structure\n\n" + _list_or_missing(sections_found, "No clear headings or section labels were detected in extracted text."),
    ]
    if condition_lines:
        sections.append("## Prompt Condition Checklist\n\n" + "\n".join(condition_lines))
    sections.append("## Technical Concepts and Methodology Signals\n\n" + _list_or_missing(methodology_terms + signals.get("models_or_algorithms", []), "No explicit methodology, model, algorithm, dataset, evaluation, or result terms were detected."))
    sections.append("## Tables, Results, and Evidence\n\n" + _list_or_missing(tables, "No tabular evidence or clearly delimited result rows were detected."))
    sections.append(
        "## Limitations and Risks\n\n"
        "- Extracted text may omit formatting, diagrams, equations, or scanned-image content.\n"
        "- Claims and results should be treated as document evidence only when visible in the extracted content.\n"
        "- If the source file is scanned or image-heavy, OCR may be required for complete analysis."
    )
    sections.append(
        "## Recommended Improvements\n\n"
        "- Add explicit methodology and result labels where missing.\n"
        "- Include captions for figures and tables.\n"
        "- Add references, reproducibility details, and implementation evidence where relevant."
    )
    sections.append(f"## Final Summary\n\nThe file `{rel}` was analyzed as a document/presentation. The report identifies structure, technical concepts, visible evidence, limitations, and recommended improvements without reproducing the document content verbatim.")
    return "\n\n".join(sections)


def _build_spreadsheet_report(settings: Settings, path: Path, instructions: str = "") -> str:
    extracted = extract_file_content(settings, path)
    rel = extracted["path"]
    content = extracted["content"]
    focus = instructions.strip() or "Create a technical report from the spreadsheet."
    sheets = re.findall(r"--- Sheet: (.*?) ---", content)
    rows = [line for line in content.splitlines() if line and not line.startswith("--- Sheet:")]
    headers = rows[0] if rows else "Not explicitly available in the provided file."
    numeric_cells = len(re.findall(r"(?<![A-Za-z])-?\d+(?:\.\d+)?", content))
    sections = [
        "# Technical Summary Report",
        f"## Overview\n\nThe spreadsheet contains {len(sheets)} detected sheet(s) and {len(rows)} sampled visible row(s). The visible first row or header is `{headers}`. The workbook presents visible sheet structure, available columns, sampled records, and basic data signals from the extracted spreadsheet content without copying the full dataset.",
        f"## File Type and Purpose\n\n- Path: `{rel}`\n- File type: `{extracted['suffix']}`\n- Size: {extracted['size']} bytes\n- Requested focus: {focus}",
        "## Workbook or Table Structure\n\n" + _list_or_missing(sheets, "No explicit sheet names were detected."),
        f"## Visible Columns or First Row\n\n{headers}",
        f"## Data Signals\n\n- Visible sampled rows: {len(rows)}\n- Numeric-looking values in extracted sample: {numeric_cells}",
        "## Limitations and Recommended Improvements\n\n- Spreadsheet extraction samples only visible rows from each sheet.\n- Full profiling should compute missing values, distributions, data types, duplicates, outliers, and target-column suitability by executing a profiling script.",
        f"## Final Summary\n\nThe file `{rel}` was summarized as a spreadsheet using visible sheet and row evidence."
    ]
    return "\n\n".join(sections)


_dispatch_previous_build_technical_report = build_technical_report


def build_technical_report(settings: Settings, path: Path, instructions: str = "") -> str:
    path = workspace.ensure_inside_root(settings, path)
    suffix = path.suffix.lower()
    if suffix == ".ipynb":
        return _build_notebook_technical_report(settings, path, instructions)
    content, kind = _extract_generic_content(path, REPORT_MAX_CHARS)
    if kind == "spreadsheet":
        return _build_spreadsheet_report(settings, path, instructions)
    if kind in _DOCUMENT_KINDS:
        return _build_document_report(settings, path, instructions)
    if kind in _SOURCE_CODE_KINDS or suffix in CODE_AND_TEXT_SUFFIXES:
        return _build_source_code_report(settings, path, instructions)
    return _dispatch_previous_build_technical_report(settings, path, instructions)

# Override metric extraction with stricter label/value parsing to avoid IDs or dataframe hashes.

def _extract_metrics(output_text: str) -> list[dict[str, str]]:
    metrics = []
    if not output_text:
        return metrics
    metric_names = [
        "mean absolute error", "mean squared error", "root mean squared error", "r-squared",
        "validation loss", "val_loss", "accuracy", "precision", "recall", "rmse", "mae", "mse",
        "r^2", "r2", "f1 score", "f1", "loss", "score",
    ]
    name_pattern = "|".join(re.escape(name) for name in metric_names)
    value_pattern = r"[-+]?\d+(?:\.\d+)?(?:[eE][-+]?\d+)?%?"
    for line in output_text.splitlines():
        line_clean = re.sub(r"\s+", " ", line).strip()
        if not line_clean or len(line_clean) > 260:
            continue
        if not re.search(r"[:=]|\b(?:is|of)\b", line_clean, flags=re.IGNORECASE):
            continue
        for match in re.finditer(rf"(?<![A-Za-z0-9_])({name_pattern})(?![A-Za-z0-9_])\s*(?:score)?\s*(?:[:=\-]|\bis\b|\bof\b)?\s*({value_pattern})(?![A-Za-z0-9_])", line_clean, flags=re.IGNORECASE):
            raw_name = match.group(1).lower()
            metric = _METRIC_ALIASES.get(raw_name, raw_name.upper())
            value = match.group(2)
            entry = {"metric": metric, "value": value, "line": line_clean[:220]}
            if entry not in metrics:
                metrics.append(entry)
    return metrics

# Override metric extraction again to support labels like Accuracy (Method Name): 0.91.

def _extract_metrics(output_text: str) -> list[dict[str, str]]:
    metrics = []
    if not output_text:
        return metrics
    metric_names = [
        "explained variance ratio", "mean absolute error", "mean squared error", "root mean squared error",
        "r-squared", "validation loss", "val_loss", "accuracy", "precision", "recall", "rmse",
        "mae", "mse", "r^2", "r2", "f1 score", "f1", "loss", "score",
    ]
    name_pattern = "|".join(re.escape(name) for name in metric_names)
    scalar_value = r"[-+]?\d+(?:\.\d+)?(?:[eE][-+]?\d+)?%?"
    list_value = r"\[[^\]]{1,180}\]"
    value_pattern = rf"(?:{list_value}|{scalar_value})"
    for raw_line in output_text.replace("\\n", "\n").splitlines():
        line_clean = re.sub(r"\s+", " ", raw_line).strip()
        if not line_clean or len(line_clean) > 320:
            continue
        if not re.search(r"[:=]", line_clean):
            continue
        pattern = rf"(?<![A-Za-z0-9_])({name_pattern})(?![A-Za-z0-9_])\s*(?:\(([^\)]{{1,120}})\))?\s*[:=]\s*({value_pattern})(?![A-Za-z0-9_])"
        for match in re.finditer(pattern, line_clean, flags=re.IGNORECASE):
            raw_name = match.group(1).lower()
            metric = _METRIC_ALIASES.get(raw_name, raw_name.title())
            label = (match.group(2) or "").strip()
            value = match.group(3).strip()
            entry = {"metric": metric, "value": value, "line": line_clean[:260]}
            if label:
                entry["label"] = label
            if entry not in metrics:
                metrics.append(entry)
    return metrics

# Chat-oriented summaries are separate from file-writing reports.

def build_chat_summary(settings: Settings, path: Path, instructions: str = "") -> str:
    path = workspace.ensure_inside_root(settings, path)
    if path.suffix.lower() == ".ipynb":
        report = _build_notebook_technical_report(settings, path, instructions or "Create a concise technical summary for chat.", include_figures=False)
    else:
        report = build_technical_report(settings, path, instructions or "Create a concise technical summary for chat.")
    lines = []
    skip_sections = {"Limitations and Risks", "Recommended Improvements", "Final Summary"}
    current = ""
    kept = 0
    for line in report.splitlines():
        heading = re.match(r"^##\s+(.+)$", line.strip())
        if heading:
            current = heading.group(1).strip()
            if current in skip_sections:
                continue
            kept += 1
            if kept > 5:
                continue
            lines.append(line)
            continue
        if current in skip_sections or kept > 5:
            continue
        if line.startswith("!["):
            continue
        if line.startswith("# "):
            continue
        lines.append(line)
    summary = "\n".join(lines).strip()
    summary = re.sub(r"\n{3,}", "\n\n", summary)
    if len(summary) > 5500:
        summary = summary[:5400].rstrip() + "\n\n[summary truncated for Discord]"
    return summary or "No readable summary could be generated from the selected file."
